"""
프로젝트 브리프용 업로드 파일에서 평문 추출 (Word/Excel/PPT/PDF 등).
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import Final

MAX_BRIEF_FILE_BYTES: Final[int] = 15 * 1024 * 1024

# 브라우저에서 UTF-8로 읽기 적합한 확장자 (클라이언트 우선 처리 가능)
PLAIN_EXTENSIONS: Final[frozenset[str]] = frozenset(
    {
        ".txt",
        ".md",
        ".markdown",
        ".csv",
        ".tsv",
        ".json",
        ".jsonl",
        ".log",
        ".xml",
        ".html",
        ".htm",
        ".xhtml",
        ".yaml",
        ".yml",
        ".css",
        ".scss",
        ".less",
        ".js",
        ".jsx",
        ".mjs",
        ".cjs",
        ".ts",
        ".tsx",
        ".vue",
        ".svelte",
        ".py",
        ".pyw",
        ".rb",
        ".go",
        ".rs",
        ".java",
        ".kt",
        ".kts",
        ".c",
        ".h",
        ".cpp",
        ".hpp",
        ".cc",
        ".cxx",
        ".cs",
        ".sql",
        ".sh",
        ".bash",
        ".zsh",
        ".ps1",
        ".env",
        ".ini",
        ".cfg",
        ".toml",
        ".properties",
        ".gitignore",
        ".dockerignore",
        ".editorconfig",
        ".svg",
    }
)


def decode_plain_bytes(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr", "shift_jis", "gb2312", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text() or ""
        parts.append(t)
    return "\n\n".join(parts).strip()


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts).strip()


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            line = "\t".join("" if c is None else str(c) for c in row)
            if line.strip():
                parts.append(line)
    wb.close()
    return "\n".join(parts).strip()


def _extract_xls(data: bytes) -> str:
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    parts: list[str] = []
    for s in book.sheets():
        parts.append(f"## {s.name}")
        for r in range(s.nrows):
            row = s.row_values(r)
            line = "\t".join(str(c) for c in row)
            if line.strip():
                parts.append(line)
    return "\n".join(parts).strip()


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
    return "\n\n".join(parts).strip()


def _extract_ods(data: bytes) -> str:
    import pandas as pd

    sheets = pd.read_excel(io.BytesIO(data), sheet_name=None, engine="odf")
    parts: list[str] = []
    for name, df in sheets.items():
        parts.append(f"## {name}")
        parts.append(df.to_csv(index=False, sep="\t"))
    return "\n\n".join(parts).strip()


def _extract_odt(data: bytes) -> str:
    from odf import teletype
    from odf.opendocument import load
    from odf.text import P

    doc = load(io.BytesIO(data))
    buf: list[str] = []
    for p in doc.getElementsByType(P):
        buf.append(teletype.extractText(p))
    return "\n".join(buf).strip()


def _extract_rtf(data: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    raw = data.decode("latin-1", errors="replace")
    return rtf_to_text(raw).strip()


def _mostly_printable(s: str, threshold: float = 0.82) -> bool:
    if not s:
        return True
    ok = sum(1 for c in s if c.isprintable() or c in "\r\n\t")
    return ok / len(s) >= threshold


def extract_text_from_bytes(filename: str, data: bytes) -> str:
    if len(data) > MAX_BRIEF_FILE_BYTES:
        raise ValueError("파일이 너무 큽니다 (최대 15MB).")

    ext = Path(filename or "upload").suffix.lower()

    if ext in PLAIN_EXTENSIONS:
        return decode_plain_bytes(data)

    if ext == ".pdf":
        return _extract_pdf(data)

    if ext == ".docx":
        return _extract_docx(data)

    if ext == ".doc":
        raise ValueError(
            "구형 Word(.doc)는 지원하지 않습니다. Word에서 .docx로 저장한 뒤 다시 올려 주세요."
        )

    if ext in (".xlsx", ".xlsm", ".xltx", ".xltm"):
        return _extract_xlsx(data)

    if ext == ".xls":
        return _extract_xls(data)

    if ext == ".pptx":
        return _extract_pptx(data)

    if ext == ".ppt":
        raise ValueError(
            "구형 PowerPoint(.ppt)는 지원하지 않습니다. .pptx로 저장한 뒤 다시 올려 주세요."
        )

    if ext == ".ods":
        return _extract_ods(data)

    if ext == ".odt":
        return _extract_odt(data)

    if ext == ".rtf":
        return _extract_rtf(data)

    # 확장자 없음·알 수 없음: UTF-8 등으로 읽어 보고, 바이너리 같으면 거절
    text = decode_plain_bytes(data)
    if _mostly_printable(text) and "\ufffd" not in text[:2000]:
        return text

    raise ValueError(
        f"이 파일 형식({ext or '확장자 없음'})에서 텍스트를 자동 추출할 수 없습니다. "
        "PDF, Word(.docx), Excel, PowerPoint(.pptx), 한글/메모장 텍스트 등으로 저장해 보세요."
    )
